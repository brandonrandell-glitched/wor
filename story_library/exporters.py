"""Export proposals to Word (.docx) and PowerPoint (.pptx)."""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt


def _strip_md_bold(text: str) -> str:
    return re.sub(r"\*\*(.+?)\*\*", r"\1", text)


def _add_bullets_docx(doc: Document, body: str) -> None:
    for line in body.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("- "):
            doc.add_paragraph(_strip_md_bold(line[2:]), style="List Bullet")
        else:
            doc.add_paragraph(_strip_md_bold(line))


def export_docx(
    account: str,
    sections: list[tuple[str, str]],
    output_path: Path,
    proposal_title: str,
) -> Path:
    doc = Document()
    doc.add_heading(f"{proposal_title}: {account}", 0)

    for title, body in sections:
        if title == "Document Metadata":
            continue
        doc.add_heading(title, level=1)
        if "\n-" in body or body.strip().startswith("- "):
            _add_bullets_docx(doc, body)
        else:
            for paragraph in body.split("\n"):
                if paragraph.strip():
                    doc.add_paragraph(_strip_md_bold(paragraph.strip()))

    doc.save(output_path)
    return output_path


def export_pptx(
    account: str,
    sections: list[tuple[str, str]],
    output_path: Path,
    proposal_title: str,
) -> Path:
    prs = Presentation()

    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = f"{proposal_title}: {account}"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = account

    bullet_layout = prs.slide_layouts[1]
    for title, body in sections:
        if title == "Document Metadata":
            continue
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = title
        tf = slide.shapes.placeholders[1].text_frame
        tf.clear()

        bullets = [
            _strip_md_bold(line.strip().lstrip("- "))
            for line in body.split("\n")
            if line.strip()
        ][:6]
        if not bullets:
            bullets = [_strip_md_bold(body.strip())]

        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

    prs.save(output_path)
    return output_path
